#!/usr/bin/env python3
import os
import sys
import argparse
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from multiprocessing import Pool, cpu_count
from tqdm import tqdm
import logging

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='[%(levelname)s] %(message)s'
)

def pdf_to_pptx(pdf_path, out_dir="output", dpi=150):
    """Convert a single PDF into a PPTX using MuPDF + python-pptx."""
    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        logging.error(f"Cannot open PDF: {pdf_path} ({e})")
        return False

    os.makedirs(out_dir, exist_ok=True)
    base = os.path.splitext(os.path.basename(pdf_path))[0]
    pptx_path = os.path.join(out_dir, base + ".pptx")

    prs = Presentation()
    # Detect first page aspect ratio
    first_page = doc.load_page(0)
    w, h = first_page.rect.width, first_page.rect.height
    pdf_ratio = w / h

    BASE_HEIGHT = 7.5  # inches
    prs.slide_height = Inches(BASE_HEIGHT)
    prs.slide_width = Inches(BASE_HEIGHT * pdf_ratio)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    logging.info(f"Processing PDF: {pdf_path} ({len(doc)} pages)")

    for page_number in range(len(doc)):
        page = doc.load_page(page_number)
        pix = page.get_pixmap(dpi=dpi)
        img_bytes = pix.tobytes("png")

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        temp_img_path = f"{base}_temp.png"
        with open(temp_img_path, "wb") as f:
            f.write(img_bytes)

        # Insert image while keeping aspect ratio
        slide.shapes.add_picture(temp_img_path, 0, 0, width=slide_width)
        os.remove(temp_img_path)

    prs.save(pptx_path)
    logging.info(f"Saved: {pptx_path}")
    return True

def find_pdfs(root):
    """Recursively find all PDF files under root."""
    pdfs = []
    for dirpath, _, filenames in os.walk(root):
        for f in filenames:
            if f.lower().endswith(".pdf"):
                pdfs.append(os.path.join(dirpath, f))
    return pdfs

def process_pdf(args):
    """Wrapper for multiprocessing pool."""
    return pdf_to_pptx(*args)

def main():
    parser = argparse.ArgumentParser(description="Convert PDF(s) to PPTX preserving aspect ratio.")
    parser.add_argument("pdf", nargs="?", help="PDF file to convert (optional).")
    parser.add_argument("--dpi", type=int, default=150, help="DPI for image rendering (default=150)")
    parser.add_argument("--out", default="output", help="Output folder (default=output/)")
    parser.add_argument("--workers", type=int, default=cpu_count(), help="Number of parallel workers (default=CPU cores)")
    args = parser.parse_args()

    # Determine list of PDFs to process
    if args.pdf:
        if not os.path.isfile(args.pdf):
            logging.error(f"File not found: {args.pdf}")
            sys.exit(1)
        pdf_list = [args.pdf]
    else:
        logging.info("No PDF specified — scanning current directory recursively...")
        pdf_list = find_pdfs(os.getcwd())
        if not pdf_list:
            logging.info("No PDFs found.")
            return
        logging.info(f"Found {len(pdf_list)} PDF(s).")

    # Prepare arguments for multiprocessing
    pool_args = [(pdf, args.out, args.dpi) for pdf in pdf_list]

    # Run in parallel with progress bar
    with Pool(processes=args.workers) as pool:
        list(tqdm(pool.imap_unordered(process_pdf, pool_args), total=len(pool_args), desc="Converting"))

if __name__ == "__main__":
    main()
